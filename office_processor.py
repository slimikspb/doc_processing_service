"""
Office document processor for Excel and PowerPoint files
Handles text extraction from XLS, XLSX, PPT, PPTX formats
"""
import os
import logging
from typing import Dict, Any, Optional, List
from pathlib import Path

logger = logging.getLogger(__name__)

class OfficeDocumentProcessor:
    """Processor for Microsoft Office documents"""
    
    def __init__(self):
        self.supported_formats = {
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls,
            '.pptx': self._process_pptx,
            '.ppt': self._process_ppt
        }
    
    def can_process(self, file_path: str) -> bool:
        """Check if file format is supported by this processor"""
        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.supported_formats
    
    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """Extract text from Office document"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if not self.can_process(file_path):
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            processor_func = self.supported_formats[file_ext]
            result = processor_func(file_path)
            
            return {
                'text': result['text'],
                'metadata': result.get('metadata', {}),
                'format': file_ext,
                'processor': 'office_processor'
            }
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            raise
    
    def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Process Excel .xlsx files"""
        try:
            import openpyxl
            from openpyxl.utils.exceptions import InvalidFileException
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            all_text = []
            sheet_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                # Extract text from all cells
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    
                    if row_text:  # Only add non-empty rows
                        sheet_text.append(' | '.join(row_text))
                
                if sheet_text:
                    sheet_content = '\n'.join(sheet_text)
                    all_text.append(f"Sheet: {sheet_name}\n{sheet_content}")
                    sheet_data[sheet_name] = {
                        'rows': len(sheet_text),
                        'max_column': sheet.max_column,
                        'max_row': sheet.max_row
                    }
            
            return {
                'text': '\n\n'.join(all_text),
                'metadata': {
                    'sheets': list(workbook.sheetnames),
                    'sheet_count': len(workbook.sheetnames),
                    'sheet_data': sheet_data
                }
            }
            
        except ImportError:
            raise Exception("openpyxl library not available for .xlsx processing")
        except InvalidFileException as e:
            raise Exception(f"Invalid Excel file: {str(e)}")
        except Exception as e:
            raise Exception(f"Error processing XLSX file: {str(e)}")
    
    def _process_xls(self, file_path: str) -> Dict[str, Any]:
        """Process Excel .xls files"""
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(file_path)
            all_text = []
            sheet_data = {}
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                sheet_name = workbook.sheet_names()[sheet_idx]
                sheet_text = []
                
                # Extract text from all cells
                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        if cell_value:  # Skip empty cells
                            row_text.append(str(cell_value))
                    
                    if row_text:  # Only add non-empty rows
                        sheet_text.append(' | '.join(row_text))
                
                if sheet_text:
                    sheet_content = '\n'.join(sheet_text)
                    all_text.append(f"Sheet: {sheet_name}\n{sheet_content}")
                    sheet_data[sheet_name] = {
                        'rows': sheet.nrows,
                        'columns': sheet.ncols
                    }
            
            return {
                'text': '\n\n'.join(all_text),
                'metadata': {
                    'sheets': workbook.sheet_names(),
                    'sheet_count': workbook.nsheets,
                    'sheet_data': sheet_data
                }
            }
            
        except ImportError:
            raise Exception("xlrd library not available for .xls processing")
        except Exception as e:
            raise Exception(f"Error processing XLS file: {str(e)}")
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint .pptx files"""
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            all_text = []
            slide_data = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Handle tables
                    if hasattr(shape, 'table'):
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_text.append(table_text)
                
                if slide_text:
                    slide_content = '\n'.join(slide_text)
                    all_text.append(f"Slide {slide_idx}:\n{slide_content}")
                    slide_data.append({
                        'slide_number': slide_idx,
                        'text_shapes': len([s for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]),
                        'total_shapes': len(slide.shapes)
                    })
            
            return {
                'text': '\n\n'.join(all_text),
                'metadata': {
                    'slide_count': len(presentation.slides),
                    'slide_data': slide_data
                }
            }
            
        except ImportError:
            raise Exception("python-pptx library not available for .pptx processing")
        except Exception as e:
            raise Exception(f"Error processing PPTX file: {str(e)}")
    
    def _process_ppt(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint .ppt files using oletools"""
        try:
            from oletools.olevba import VBA_Parser
            from oletools import olefile
            
            # For legacy .ppt files, we'll extract what we can
            # This is more limited than .pptx processing
            
            if olefile.isOleFile(file_path):
                # Try to extract basic text content
                with open(file_path, 'rb') as f:
                    # This is a simplified extraction - .ppt format is complex
                    # In production, you might want to use a more sophisticated approach
                    content = f.read()
                    
                    # Look for readable text (basic approach)
                    text_content = self._extract_readable_text(content)
                    
                    return {
                        'text': text_content,
                        'metadata': {
                            'format': 'legacy_ppt',
                            'extraction_method': 'basic_text_search',
                            'note': 'Legacy PPT format has limited text extraction capabilities'
                        }
                    }
            else:
                raise Exception("File is not a valid OLE file")
                
        except ImportError:
            raise Exception("oletools library not available for .ppt processing")
        except Exception as e:
            raise Exception(f"Error processing PPT file: {str(e)}")
    
    def _extract_table_text(self, table) -> str:
        """Extract text from PowerPoint table"""
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                table_text.append(' | '.join(row_text))
        return '\n'.join(table_text)
    
    def _extract_readable_text(self, content: bytes) -> str:
        """Extract readable text from binary content (basic approach)"""
        try:
            # Simple approach: find printable ASCII text
            import string
            
            text_chunks = []
            current_chunk = []
            
            for byte in content:
                char = chr(byte) if byte < 128 else None
                
                if char and (char in string.printable and char not in '\x0b\x0c'):
                    current_chunk.append(char)
                else:
                    if len(current_chunk) > 10:  # Only keep chunks with reasonable length
                        chunk_text = ''.join(current_chunk).strip()
                        if chunk_text and not chunk_text.isspace():
                            text_chunks.append(chunk_text)
                    current_chunk = []
            
            # Final chunk
            if len(current_chunk) > 10:
                chunk_text = ''.join(current_chunk).strip()
                if chunk_text and not chunk_text.isspace():
                    text_chunks.append(chunk_text)
            
            return '\n'.join(text_chunks)
            
        except Exception as e:
            logger.warning(f"Error in basic text extraction: {e}")
            return "Text extraction from legacy PPT format failed"

# Global instance
office_processor = OfficeDocumentProcessor()