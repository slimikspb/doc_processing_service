"""
Reliable document text extractor using stable libraries instead of textract.
Supports PDF, DOCX, XLSX, PPTX, RTF, and TXT files.
"""

import os
import io
import logging
from pathlib import Path
from typing import Dict, Any, Optional

# PDF processing
try:
    import fitz  # PyMuPDF
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Office document processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    from xlrd import open_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Text processing
import chardet
from PIL import Image
# Raster detection
try:
    from pdf_raster_detector import detect_pdf_raster_images
    RASTER_DETECTION_AVAILABLE = True
except ImportError:
    RASTER_DETECTION_AVAILABLE = False
logger = logging.getLogger(__name__)

class ReliableDocumentExtractor:
    """
    Document text extractor using reliable libraries instead of textract.
    """
    
    def __init__(self):
        self.supported_formats = self._get_supported_formats()
        logger.info(f"Initialized with support for: {', '.join(self.supported_formats)}")
    
    def _get_supported_formats(self) -> list:
        """Get list of supported file formats based on available libraries."""
        formats = ['txt', 'rtf']  # Always supported
        
        if PDF_AVAILABLE:
            formats.extend(['pdf'])
        if DOCX_AVAILABLE:
            formats.extend(['docx', 'doc'])
        if EXCEL_AVAILABLE:
            formats.extend(['xlsx', 'xls'])
        if PPTX_AVAILABLE:
            formats.extend(['pptx'])
            
        return formats
    
    def extract_text(self, file_path: str, ocr_enabled: bool = False) -> Dict[str, Any]:
        """
        Extract text from a document file.
        
        Args:
            file_path: Path to the document file
            ocr_enabled: Whether to perform OCR on images in PDFs
            
        Returns:
            Dict with 'text' and 'metadata' keys
        """
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_ext = Path(file_path).suffix.lower().lstrip('.')
            file_size = os.path.getsize(file_path)
            
            logger.info(f"Extracting text from {file_ext.upper()} file: {file_path} ({file_size} bytes), OCR: {ocr_enabled}")
            
            # Route to appropriate extractor
            if file_ext == 'pdf':
                text = self._extract_pdf(file_path, ocr_enabled=ocr_enabled)
            elif file_ext in ['docx', 'doc']:
                text = self._extract_docx(file_path)
            elif file_ext in ['xlsx', 'xls']:
                text = self._extract_excel(file_path)
            elif file_ext == 'pptx':
                text = self._extract_pptx(file_path)
            elif file_ext in ['txt', 'rtf']:
                text = self._extract_text_file(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            return {
                'text': text.strip(),
                'metadata': {
                    'file_type': file_ext,
                    'file_size': file_size,
                    'text_length': len(text),
                    'extractor': 'reliable_extractor',
                    'ocr_enabled': ocr_enabled if file_ext == 'pdf' else False
                }
            }
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    def _extract_pdf(self, file_path: str, ocr_enabled: bool = False) -> str:
        """Extract text from PDF using PyMuPDF with pdfplumber fallback and optional OCR."""
        if not PDF_AVAILABLE:
            raise ImportError("PDF processing libraries not available")
        
        try:
            # Try PyMuPDF first (faster) - extract page by page
            doc = fitz.open(file_path)
            pages_text = []
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                pages_text.append({
                    'page_number': page_num + 1,
                    'text': page.get_text()
                })
            
            doc.close()
            
            # If PyMuPDF didn't extract much text, try pdfplumber
            total_text = '\n'.join([p['text'] for p in pages_text])
            if len(total_text.strip()) < 100:
                logger.info("PyMuPDF extracted little text, trying pdfplumber")
                with pdfplumber.open(file_path) as pdf:
                    pages_text = []
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        pages_text.append({
                            'page_number': i + 1,
                            'text': page_text if page_text else ''
                        })
            
            # Perform OCR if enabled and images are present
            if ocr_enabled:
                text = self._enrich_with_ocr(file_path, pages_text)
            else:
                # Just combine all pages
                text = '\n'.join([p['text'] for p in pages_text])
            
            return text
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            return f"Error extracting PDF: {str(e)}"
    
    def _enrich_with_ocr(self, file_path: str, pages_text: list) -> str:
        """
        Enrich PDF text with OCR from embedded images.
        
        Args:
            file_path: Path to PDF file
            pages_text: List of dicts with page_number and text for each page
            
        Returns:
            Text enriched with OCR results inline
        """
        try:
            from image_extractor import extract_images_from_pdf, cleanup_images, is_image_extraction_available
            from ocr_processor import process_images, is_ocr_available
            
            if not is_image_extraction_available():
                logger.warning("Image extraction not available, skipping OCR")
                return '\n'.join([p['text'] for p in pages_text])
            
            if not is_ocr_available():
                logger.warning("OCR not available, skipping OCR enrichment")
                return '\n'.join([p['text'] for p in pages_text])
            
            # Extract images from PDF
            images = extract_images_from_pdf(file_path)
            
            if not images:
                logger.info("No images found in PDF, skipping OCR")
                return '\n'.join([p['text'] for p in pages_text])
            
            logger.info(f"Found {len(images)} images in PDF, performing OCR")
            
            # Perform OCR on all images
            try:
                images_with_ocr = process_images(images)
                
                # Build enriched text with OCR inline
                enriched_text = self._combine_text_with_ocr(pages_text, images_with_ocr)
                
                return enriched_text
                
            finally:
                # Always cleanup temporary image files
                cleanup_images(images)
            
        except Exception as e:
            logger.error(f"OCR enrichment failed: {str(e)}")
            # Return original text if OCR fails
            return '\n'.join([p['text'] for p in pages_text])
    
    def _combine_text_with_ocr(self, pages_text: list, images_with_ocr: list) -> str:
        """
        Combine original PDF text with OCR results from images inline.
        
        Args:
            pages_text: List of dicts with page_number and text for each page
            images_with_ocr: List of image metadata with OCR results
            
        Returns:
            Combined text with OCR inline at page level
        """
        # Group images by page
        images_by_page = {}
        for img in images_with_ocr:
            page_num = img.get('page_number', 0)
            if page_num not in images_by_page:
                images_by_page[page_num] = []
            images_by_page[page_num].append(img)
        
        # Build combined text page by page
        combined_parts = []
        
        for page_info in pages_text:
            page_num = page_info['page_number']
            page_text = page_info['text']
            
            # Add the page's original text
            if page_text.strip():
                combined_parts.append(page_text)
            
            # Add OCR text from images on this page (inline)
            if page_num in images_by_page:
                for img in images_by_page[page_num]:
                    ocr_text = img.get('ocr_text', '').strip()
                    if ocr_text:
                        width = img.get('width', 0)
                        height = img.get('height', 0)
                        
                        combined_parts.append(f"\n[IMAGE ON PAGE {page_num}: {width}x{height}px]")
                        combined_parts.append(ocr_text)
                        combined_parts.append("[END IMAGE]\n")
        
        return '\n'.join(combined_parts)
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX/DOC files."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx library not available")
        
        try:
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraph text
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract table text
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            return f"Error extracting DOCX: {str(e)}"
    
    def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel files (XLSX/XLS)."""
        if not EXCEL_AVAILABLE:
            raise ImportError("Excel processing libraries not available")
        
        try:
            file_ext = Path(file_path).suffix.lower()
            text_parts = []
            
            if file_ext == '.xlsx':
                # Use openpyxl for .xlsx files
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_parts.append(f"Sheet: {sheet_name}")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None and str(cell).strip():
                                row_text.append(str(cell).strip())
                        if row_text:
                            text_parts.append(' | '.join(row_text))
                
                workbook.close()
                
            else:
                # Use xlrd for .xls files
                workbook = open_workbook(file_path)
                
                for sheet_idx in range(workbook.nsheets):
                    sheet = workbook.sheet_by_index(sheet_idx)
                    text_parts.append(f"Sheet: {sheet.name}")
                    
                    for row_idx in range(sheet.nrows):
                        row_text = []
                        for col_idx in range(sheet.ncols):
                            cell = sheet.cell(row_idx, col_idx)
                            if cell.value and str(cell.value).strip():
                                row_text.append(str(cell.value).strip())
                        if row_text:
                            text_parts.append(' | '.join(row_text))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Excel extraction failed: {str(e)}")
            return f"Error extracting Excel: {str(e)}"
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx library not available")
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {slide_idx}:")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                text_parts.append("")  # Empty line between slides
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            return f"Error extracting PPTX: {str(e)}"
    
    def detect_raster_images(self, file_path: str, settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Detect raster images in PDF files.
        
        Args:
            file_path: Path to the PDF file
            settings: Optional configuration settings for detection
            
        Returns:
            Dict with raster image detection results
        """
        if not RASTER_DETECTION_AVAILABLE:
            raise ImportError("Raster detection not available - missing dependencies")
        
        file_ext = Path(file_path).suffix.lower().lstrip(".")
        if file_ext != "pdf":
            raise ValueError("Raster detection only supported for PDF files")
        
        try:
            return detect_pdf_raster_images(file_path, settings)
        except Exception as e:
            logger.error(f"Error detecting raster images in {file_path}: {str(e)}")
            raise Exception(f"Raster detection failed: {str(e)}")
    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text files with encoding detection."""
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result['encoding'] or 'utf-8'
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                return f.read()
                
        except Exception as e:
            logger.error(f"Text file extraction failed: {str(e)}")
            # Fallback to UTF-8 with error replacement
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            except Exception as fallback_error:
                return f"Error reading text file: {str(fallback_error)}"

# Global instance
reliable_extractor = ReliableDocumentExtractor()

def extract_document_text(file_path: str, ocr_enabled: bool = False) -> Dict[str, Any]:
    """
    Convenience function to extract text from a document.
    
    Args:
        file_path: Path to the document file
        ocr_enabled: Whether to perform OCR on images in PDFs
        
    Returns:
        Dict with extracted text and metadata
    """
    return reliable_extractor.extract_text(file_path, ocr_enabled=ocr_enabled)

def detect_pdf_raster(file_path: str, settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """
    Convenience function to detect raster images in PDF.
    
    Args:
        file_path: Path to the PDF file
        settings: Optional configuration settings
        
    Returns:
        Dict containing raster image detection results
    """
    return reliable_extractor.detect_raster_images(file_path, settings)
def get_supported_formats() -> list:
    """Get list of supported document formats."""
    return reliable_extractor.supported_formats