#!/usr/bin/env python3
"""
Test suite for Office document processing (Excel and PowerPoint)
"""
import requests
import argparse
import json
import sys
import tempfile
import os
from pathlib import Path

def create_test_excel_file(filepath: str, format_type: str = 'xlsx'):
    """Create a test Excel file for testing"""
    try:
        if format_type == 'xlsx':
            import openpyxl
            
            # Create a new workbook with test data
            wb = openpyxl.Workbook()
            
            # Sheet 1: Basic data
            ws1 = wb.active
            ws1.title = "Sales Data"
            ws1['A1'] = "Product"
            ws1['B1'] = "Quantity"
            ws1['C1'] = "Price"
            ws1['A2'] = "Widget A"
            ws1['B2'] = 100
            ws1['C2'] = 19.99
            ws1['A3'] = "Widget B"
            ws1['B3'] = 50
            ws1['C3'] = 29.99
            
            # Sheet 2: Text content
            ws2 = wb.create_sheet("Documentation")
            ws2['A1'] = "Product Documentation"
            ws2['A2'] = "This spreadsheet contains sales data for Q1 2024"
            ws2['A3'] = "Key metrics include product names, quantities sold, and unit prices"
            ws2['A4'] = "Total products tracked: 2"
            
            wb.save(filepath)
            return True
            
        elif format_type == 'xls':
            import xlwt
            
            # Create a new workbook for XLS format
            wb = xlwt.Workbook()
            
            # Sheet 1: Basic data
            ws1 = wb.add_sheet("Sales Data")
            ws1.write(0, 0, "Product")
            ws1.write(0, 1, "Quantity")
            ws1.write(0, 2, "Price")
            ws1.write(1, 0, "Widget A")
            ws1.write(1, 1, 100)
            ws1.write(1, 2, 19.99)
            ws1.write(2, 0, "Widget B")
            ws1.write(2, 1, 50)
            ws1.write(2, 2, 29.99)
            
            # Sheet 2: Text content
            ws2 = wb.add_sheet("Documentation")
            ws2.write(0, 0, "Product Documentation")
            ws2.write(1, 0, "This spreadsheet contains sales data for Q1 2024")
            ws2.write(2, 0, "Key metrics include product names, quantities sold, and unit prices")
            
            wb.save(filepath)
            return True
            
    except ImportError as e:
        print(f"⚠ Cannot create {format_type} test file: {e}")
        return False
    except Exception as e:
        print(f"✗ Error creating test {format_type} file: {e}")
        return False

def create_test_powerpoint_file(filepath: str, format_type: str = 'pptx'):
    """Create a test PowerPoint file for testing"""
    try:
        if format_type == 'pptx':
            from pptx import Presentation
            from pptx.util import Inches
            
            # Create a new presentation
            prs = Presentation()
            
            # Slide 1: Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(title_slide_layout)
            title = slide1.shapes.title
            subtitle = slide1.placeholders[1]
            title.text = "Quarterly Business Review"
            subtitle.text = "Q1 2024 Performance Summary"
            
            # Slide 2: Bullet points
            bullet_slide_layout = prs.slide_layouts[1]
            slide2 = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide2.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = "Key Achievements"
            tf = body_shape.text_frame
            tf.text = "Increased sales by 25%"
            
            p = tf.add_paragraph()
            p.text = "Launched 3 new products"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = "Expanded to 2 new markets"
            p.level = 0
            
            # Slide 3: Content slide
            blank_slide_layout = prs.slide_layouts[6]
            slide3 = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_shape = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = "Financial Summary"
            
            # Add content
            content_shape = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            content_frame = content_shape.text_frame
            content_frame.text = "Revenue: $1,250,000\\nExpenses: $980,000\\nNet Profit: $270,000\\nGrowth Rate: 15%"
            
            prs.save(filepath)
            return True
            
    except ImportError as e:
        print(f"⚠ Cannot create {format_type} test file: {e}")
        return False
    except Exception as e:
        print(f"✗ Error creating test {format_type} file: {e}")
        return False

def test_excel_processing(base_url: str, api_key: str, file_format: str = 'xlsx'):
    """Test Excel file processing"""
    print(f"Testing Excel {file_format.upper()} processing...")
    
    # Create temporary test file
    with tempfile.NamedTemporaryFile(suffix=f'.{file_format}', delete=False) as tmp_file:
        temp_path = tmp_file.name
    
    try:
        # Create test Excel file
        if not create_test_excel_file(temp_path, file_format):
            print(f"✗ Could not create test {file_format} file")
            return False
        
        headers = {'X-API-Key': api_key} if api_key else {}
        
        with open(temp_path, 'rb') as f:
            files = {'document': f}
            response = requests.post(f"{base_url}/convert", 
                                   headers=headers, 
                                   files=files, 
                                   timeout=60)
        
        if response.status_code == 200:
            data = response.json()
            print(f"✓ Excel {file_format.upper()} processing successful")
            print(f"  Filename: {data.get('filename', 'unknown')}")
            print(f"  Extraction method: {data.get('extraction_method', 'unknown')}")
            print(f"  Text length: {len(data.get('text', ''))}")
            
            # Check for expected content
            text = data.get('text', '').lower()
            if 'widget' in text and 'sales data' in text:
                print("  ✓ Expected content found in extracted text")
            else:
                print("  ⚠ Expected content not found in extracted text")
            
            # Check metadata
            metadata = data.get('metadata', {})
            if 'sheets' in metadata or 'sheet_count' in metadata:
                print(f"  ✓ Metadata extracted: {metadata}")
            
            return True
        else:
            print(f"✗ Excel {file_format.upper()} processing failed: {response.status_code} - {response.text}")
            return False
            
    except Exception as e:
        print(f"✗ Excel {file_format.upper()} test error: {e}")
        return False
    finally:
        # Clean up temp file
        try:
            os.unlink(temp_path)
        except:
            pass

def test_powerpoint_processing(base_url: str, api_key: str, file_format: str = 'pptx'):
    """Test PowerPoint file processing"""
    print(f"Testing PowerPoint {file_format.upper()} processing...")
    
    # Create temporary test file
    with tempfile.NamedTemporaryFile(suffix=f'.{file_format}', delete=False) as tmp_file:
        temp_path = tmp_file.name
    
    try:
        # Create test PowerPoint file
        if not create_test_powerpoint_file(temp_path, file_format):
            print(f"✗ Could not create test {file_format} file")
            return False
        
        headers = {'X-API-Key': api_key} if api_key else {}
        
        with open(temp_path, 'rb') as f:
            files = {'document': f}
            response = requests.post(f"{base_url}/convert", 
                                   headers=headers, 
                                   files=files, 
                                   timeout=60)
        
        if response.status_code == 200:
            data = response.json()
            print(f"✓ PowerPoint {file_format.upper()} processing successful")
            print(f"  Filename: {data.get('filename', 'unknown')}")
            print(f"  Extraction method: {data.get('extraction_method', 'unknown')}")
            print(f"  Text length: {len(data.get('text', ''))}")
            
            # Check for expected content
            text = data.get('text', '').lower()
            if 'quarterly' in text and 'business review' in text:
                print("  ✓ Expected content found in extracted text")
            else:
                print("  ⚠ Expected content not found in extracted text")
            
            # Check metadata
            metadata = data.get('metadata', {})
            if 'slide_count' in metadata:
                print(f"  ✓ Metadata extracted: {metadata}")
            
            return True
        else:
            print(f"✗ PowerPoint {file_format.upper()} processing failed: {response.status_code} - {response.text}")
            return False
            
    except Exception as e:
        print(f"✗ PowerPoint {file_format.upper()} test error: {e}")
        return False
    finally:
        # Clean up temp file
        try:
            os.unlink(temp_path)
        except:
            pass

def test_supported_formats(base_url: str):
    """Test the supported formats endpoint"""
    print("Testing supported formats endpoint...")
    
    try:
        response = requests.get(f"{base_url}/formats", timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            print("✓ Supported formats endpoint working")
            
            office_formats = data.get('office_documents', {}).get('formats', [])
            textract_formats = data.get('textract_documents', {}).get('formats', [])
            total = data.get('total_supported', 0)
            
            print(f"  Office formats: {office_formats}")
            print(f"  Textract formats: {len(textract_formats)} formats")
            print(f"  Total supported: {total}")
            
            # Check if our expected formats are listed
            expected_office = ['.xlsx', '.xls', '.pptx', '.ppt']
            missing = [fmt for fmt in expected_office if fmt not in office_formats]
            if not missing:
                print("  ✓ All expected Office formats are supported")
            else:
                print(f"  ⚠ Missing Office formats: {missing}")
            
            return True
        else:
            print(f"✗ Supported formats endpoint failed: {response.status_code}")
            return False
            
    except Exception as e:
        print(f"✗ Supported formats test error: {e}")
        return False

def run_office_tests(base_url: str, api_key: str):
    """Run all Office document tests"""
    print(f"Running Office document tests against {base_url}")
    print("=" * 60)
    
    results = []
    
    # Test supported formats endpoint
    results.append(test_supported_formats(base_url))
    
    # Test Excel formats
    results.append(test_excel_processing(base_url, api_key, 'xlsx'))
    results.append(test_excel_processing(base_url, api_key, 'xls'))
    
    # Test PowerPoint formats
    results.append(test_powerpoint_processing(base_url, api_key, 'pptx'))
    # Note: Legacy PPT testing is limited due to format complexity
    
    print("=" * 60)
    passed = sum(results)
    total = len(results)
    
    print(f"Office Document Tests: {passed}/{total} passed")
    
    if passed == total:
        print("🎉 All Office document tests passed!")
        return True
    else:
        print("⚠ Some Office document tests failed")
        return False

def main():
    parser = argparse.ArgumentParser(description='Test suite for Office document processing')
    parser.add_argument('--url', default='http://localhost:5001', 
                       help='Base URL of the service (default: http://localhost:5001)')
    parser.add_argument('--api-key', default='default_dev_key',
                       help='API key for authentication (default: default_dev_key)')
    parser.add_argument('--test', choices=['formats', 'excel', 'powerpoint', 'all'],
                       default='all', help='Specific test to run')
    
    args = parser.parse_args()
    
    if args.test == 'formats':
        success = test_supported_formats(args.url)
    elif args.test == 'excel':
        success = (test_excel_processing(args.url, args.api_key, 'xlsx') and
                  test_excel_processing(args.url, args.api_key, 'xls'))
    elif args.test == 'powerpoint':
        success = test_powerpoint_processing(args.url, args.api_key, 'pptx')
    else:  # all
        success = run_office_tests(args.url, args.api_key)
    
    sys.exit(0 if success else 1)

if __name__ == '__main__':
    main()